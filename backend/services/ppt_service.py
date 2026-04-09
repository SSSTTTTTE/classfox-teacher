"""
课程资料解析服务
================
支持 PPT (.pptx)、PDF (.pdf)、Word (.docx) 文件解析为纯文本
"""

import os
from pptx import Presentation


def parse_ppt_to_text(file_path: str) -> str:
    """
    解析 PPT 文件，提取所有幻灯片中的文本内容
    """
    prs = Presentation(file_path)
    all_text = []

    for slide_num, slide in enumerate(prs.slides, 1):
        slide_texts = []
        slide_texts.append(f"--- 第 {slide_num} 页 ---")

        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        slide_texts.append(text)

            if shape.has_table:
                table = shape.table
                for row in table.rows:
                    row_text = []
                    for cell in row.cells:
                        cell_text = cell.text.strip()
                        if cell_text:
                            row_text.append(cell_text)
                    if row_text:
                        slide_texts.append(" | ".join(row_text))

        if len(slide_texts) > 1:
            all_text.append("\n".join(slide_texts))

    return "\n\n".join(all_text)


def parse_pdf_to_text(file_path: str) -> str:
    """
    解析 PDF 文件，提取所有页面中的文本内容
    """
    from pypdf import PdfReader

    reader = PdfReader(file_path)
    all_text = []

    for i, page in enumerate(reader.pages, 1):
        text = page.extract_text()
        if text and text.strip():
            all_text.append(f"--- 第 {i} 页 ---\n{text.strip()}")

    return "\n\n".join(all_text)


def parse_docx_to_text(file_path: str) -> str:
    """
    解析 Word (.docx) 文件，提取所有段落和表格中的文本
    """
    from docx import Document

    doc = Document(file_path)
    all_text = []

    for para in doc.paragraphs:
        text = para.text.strip()
        if text:
            all_text.append(text)

    for table in doc.tables:
        for row in table.rows:
            row_text = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if row_text:
                all_text.append(" | ".join(row_text))

    return "\n".join(all_text)


def parse_material(file_path: str, filename: str) -> str:
    """
    根据文件扩展名自动选择解析器

    Args:
        file_path: 文件绝对路径
        filename: 原始文件名（用于判断扩展名）

    Returns:
        提取出的纯文本内容
    """
    ext = os.path.splitext(filename)[1].lower()

    if ext in (".pptx", ".ppt"):
        return parse_ppt_to_text(file_path)
    elif ext == ".pdf":
        return parse_pdf_to_text(file_path)
    elif ext in (".docx", ".doc"):
        return parse_docx_to_text(file_path)
    else:
        raise ValueError(f"不支持的文件格式: {ext}")
